"""Generate three FedBook-Sem presentation decks (PPTX + PDF) from a single
content source.

Outputs (next to this file):
    FedBookSem_Academic_Viva.{pptx,pdf}
    FedBookSem_Technical_Walkthrough.{pptx,pdf}
    FedBookSem_General_Overview.{pptx,pdf}

Requires: python-pptx, reportlab
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from reportlab.lib import colors as rl_colors
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import mm
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Paragraph

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

PRIMARY = "#1E2A44"   # deep navy
ACCENT = "#E9B44C"    # warm gold
INK = "#1A1A1A"       # near-black text
MUTED = "#667085"     # secondary text
BG = "#FFFFFF"        # slide background
SOFT = "#F5F6F8"      # subtle panels

SLIDE_W_IN = 13.333   # 16:9 wide
SLIDE_H_IN = 7.5


def _hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _hex_to_rl(h: str) -> rl_colors.Color:
    h = h.lstrip("#")
    return rl_colors.Color(int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255)


# ---------------------------------------------------------------------------
# Content — a single source used to render every deck.
#
# Each slide is a dict with a "kind" and kind-specific fields.
# Kinds:
#   title      { title, subtitle, footer }
#   section    { title, kicker? }
#   bullets    { title, items: [str], note? }
#   two_col    { title, left: {heading, items}, right: {heading, items} }
#   kv         { title, rows: [(k, v)] }
#   quote      { title, quote, attribution? }
# ---------------------------------------------------------------------------

# Shared building blocks -----------------------------------------------------

PROBLEM_STATEMENT = (
    "Centralised bookstores and reading platforms own the social graph, the "
    "reviews, and the annotations. Readers cannot move between services "
    "without losing their followers or their work, and cross-platform sentiment "
    "signal is invisible to the recommender."
)

SOLUTION_STATEMENT = (
    "FedBook-Sem is a federated social layer for an intelligent online "
    "bookstore. Users own an ActivityPub actor, write W3C Web Annotations on "
    "book passages, and receive recommendations that fuse semantic similarity, "
    "collaborative filtering, and privacy-preserving cross-platform reception."
)

TECH_ROWS = [
    ("Federation", "ActivityPub + WebFinger + HTTP Signatures (RSA-2048)"),
    ("Annotations", "W3C Web Annotation Data Model — TextQuote + TextPosition selectors"),
    ("Backend", "Node.js 20 · Express · JWT · bcrypt · node-forge · multer · Bull"),
    ("Frontend", "React 18 · React Router v6 · Axios · custom dark design system"),
    ("Graph DB", "Neo4j 5.15 + APOC + Graph Data Science + native vector index"),
    ("ML service", "Python 3.11 · FastAPI · sentence-transformers · LightFM · LightGBM"),
    ("Sentiment", "CardiffNLP twitter-roberta-base-sentiment-latest (RoBERTa)"),
    ("Cross-encoder", "BAAI/bge-reranker-base (~500 MB)"),
    ("Dashboard", "Streamlit — Recommender demo + Model registry"),
    ("Infrastructure", "Docker Compose · Redis 7 · PostgreSQL 16"),
]

RECOMMENDER_ROWS = [
    ("Semantic", "MiniLM-L6-v2 embeddings + Neo4j native vector index (cosine k-NN)"),
    ("Cross-encoder rerank", "BAAI/bge-reranker-base for query × candidate re-scoring"),
    ("Collaborative filtering", "LightFM WARP trained on goodbooks-10k (~6M real ratings)"),
    ("Graph traversal", "Neo4j GDS k-NN and personalised PageRank"),
    ("Cross-platform reception", "Weighted sentiment across Reddit / YouTube / Bluesky / Mastodon"),
    ("Learned rank fusion", "LightGBM LambdaRank over 6 features (sim, cf, reception, diversity, ce, freshness)"),
]

INGEST_ROWS = [
    ("Reddit", "PRAW OAuth · subreddits /r/books, /r/booksuggestions, /r/52book, /r/literature"),
    ("YouTube", "Data API v3 · top-level comments on review videos · 30-day retention"),
    ("Bluesky", "AT Protocol public searchPosts · optional app-password auth"),
    ("Mastodon", "Public hashtag timelines · #bookstodon, #booksky, #booktok, #bookreview"),
]

# ----- Academic viva deck ---------------------------------------------------

ACADEMIC_SLIDES: list[dict[str, Any]] = [
    {
        "kind": "title",
        "title": "FedBook-Sem",
        "subtitle": "A Federated Social Layer for an Intelligent Online Bookstore",
        "footer": "R26-IT-107  ·  IT22922670  ·  SLIIT B.Sc. (Hons) 2026",
    },
    {
        "kind": "section",
        "kicker": "01",
        "title": "Motivation & Research Gap",
    },
    {
        "kind": "bullets",
        "title": "The Problem",
        "items": [
            "Book platforms are silos — Goodreads, Amazon, StoryGraph each keep their own social graph.",
            "Reader-generated annotations are locked to one product and lost on migration.",
            "Recommenders ignore the wider social conversation happening on Reddit, YouTube, Bluesky and Mastodon.",
            "There is no open specification for interoperable, scholarly-grade book annotations.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Research Gap",
        "items": [
            "Federated social protocols (ActivityPub) are proven for microblogging (Mastodon), but under-used for books.",
            "W3C Web Annotation is standardised for scholarly texts, rarely applied to trade-fiction reading UIs.",
            "Cross-platform reception is rarely folded into a single recommender score — most systems use only intrinsic ratings.",
            "No existing system combines all three in one interoperable stack.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Research Questions",
        "items": [
            "RQ1 — Can ActivityPub sustain a book-focused social network with W3C annotations as first-class activities?",
            "RQ2 — Does adding privacy-preserving cross-platform reception measurably improve recommendation quality (NDCG@10, MRR) over semantic or CF alone?",
            "RQ3 — Can learned rank fusion (LightGBM LambdaRank) combine heterogeneous signals better than a fixed-weight linear ensemble?",
        ],
    },
    {
        "kind": "section",
        "kicker": "02",
        "title": "Contribution",
    },
    {
        "kind": "bullets",
        "title": "Contributions",
        "items": [
            "A federated book-social protocol layered on ActivityPub with W3C annotations as CREATE activities.",
            "A privacy-preserving ingestion pipeline that scores cross-platform sentiment without persisting raw mention text.",
            "A hybrid recommender combining semantic + CF + reception + graph signals, with an LTR fusion head.",
            "An open, reproducible dev environment (Docker Compose, seedable Neo4j, offline-cached models).",
            "A Streamlit evaluation dashboard exposing every recommender strategy side-by-side.",
        ],
    },
    {
        "kind": "section",
        "kicker": "03",
        "title": "Methodology",
    },
    {
        "kind": "bullets",
        "title": "Methodology — Design Science",
        "items": [
            "Design Science Research Methodology (Peffers et al.) — iterative build-and-evaluate cycles.",
            "Six artefacts: federation protocol, annotation model, ingestion pipeline, recommender ensemble, dashboard, evaluation harness.",
            "Each artefact justified against literature, then evaluated with a mix of technical and human metrics.",
            "Reproducibility: single command (docker compose up -d) reconstructs the entire stack from source.",
        ],
    },
    {
        "kind": "two_col",
        "title": "Evaluation Plan",
        "left": {
            "heading": "Recommender quality (offline)",
            "items": [
                "NDCG@5, NDCG@10, MRR, Recall@20",
                "Held-out split on goodbooks-10k",
                "Ablation: semantic-only vs. +reception vs. +CF vs. full LTR",
                "Coverage & catalogue diversity (Shannon entropy)",
            ],
        },
        "right": {
            "heading": "System & user (mixed)",
            "items": [
                "Federation conformance vs. Mastodon reference server",
                "Annotation round-trip (W3C validator)",
                "10-user think-aloud study, task-completion time",
                "Latency: P50 / P95 on /recommend/similar",
            ],
        },
    },
    {
        "kind": "bullets",
        "title": "Ethical & Legal Considerations",
        "items": [
            "Raw mention text from Reddit / YouTube / Bluesky / Mastodon is never persisted — only aggregated sentiment counts and external IDs for the deletion sweep.",
            "YouTube data honours the 30-day retention window (purge_expired_receptions.py).",
            "Reddit deletion sweep (reddit_deletion_sweep.py) removes counts for deleted posts.",
            "All external APIs are called with a project user-agent identifying the researcher.",
            "User study cleared by SLIIT research ethics; consent form + PIS on file.",
        ],
    },
    {
        "kind": "section",
        "kicker": "04",
        "title": "System",
    },
    {
        "kind": "kv",
        "title": "Technology Stack",
        "rows": TECH_ROWS,
    },
    {
        "kind": "quote",
        "title": "Architecture at a Glance",
        "quote": (
            "React FE  →  Express API  →  Neo4j graph, shared with FastAPI ML service.\n"
            "Streamlit dashboard talks to the ML API. A daily ingestion worker pulls "
            "book mentions from four social platforms, scores sentiment, and writes "
            ":PlatformReception nodes into the same Neo4j."
        ),
        "attribution": "See FedBookSem_Architecture.drawio for the full layered diagram.",
    },
    {
        "kind": "kv",
        "title": "Layered Architecture — Five Tiers",
        "rows": [
            ("Client", "React FE (localhost:3000)  ·  Streamlit dashboard (:8501)  ·  remote ActivityPub servers"),
            ("Application", "Express backend (:3001)  ·  FastAPI ML service (:8000)  ·  ActivityPub inbox / outbox / WebFinger"),
            ("Ingestion", "Reddit / YouTube / Bluesky / Mastodon collectors → entity resolver → sentiment scorer → aggregator"),
            ("Data", "Neo4j 5.15 (bolt :7687)  ·  Redis 7 (:6379)  ·  PostgreSQL 16 (:5432)  ·  cached model files on disk"),
            ("External", "Open Library  ·  Hardcover GraphQL  ·  HuggingFace Hub  ·  Kaggle 7k  ·  goodbooks-10k  ·  remote AP peers"),
        ],
    },
    {
        "kind": "bullets",
        "title": "Request Lifecycle — Login and Feed",
        "items": [
            "Browser POSTs /api/auth/login → CRA dev proxy forwards to Express :3001.",
            "Express queries Neo4j MATCH (:Person {username}), bcrypt.compare, jwt.sign — returns JWT.",
            "Browser stores JWT in localStorage; axios interceptor attaches it as Bearer to every subsequent call.",
            "Browser GET /api/feed → Express verifies JWT → single Cypher: MATCH (me)-[:FOLLOWS]->(f)-[:AUTHORED]->(r)-[:REVIEWS]->(b).",
            "Response JSON hydrated on the client into a virtualised feed component.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Federation Flow — Alice Follows @bob@remote.social",
        "items": [
            "Alice clicks Follow → POST /api/social/follow with targetId.",
            "Backend WebFingers acct:bob@remote.social → discovers bob's actor URL.",
            "Fetches bob's actor JSON-LD, extracts publicKey and inbox URL.",
            "Constructs an ActivityPub Follow activity, signs it with alice's private key using HTTP Signatures (RSA-2048, SHA-256).",
            "POSTs the signed activity to bob's inbox; remote server verifies signature and MERGEs the follower.",
            "Remote server later sends an Accept activity back to alice's /inbox → verified → MERGE (alice)-[:FOLLOWS]->(bob).",
        ],
    },
    {
        "kind": "bullets",
        "title": "Recommender Request Flow — /recommend/similar",
        "items": [
            "Client GET /recommend/similar?isbn=X&rerank=ltr → FastAPI.",
            "FastAPI reads :Book(embedding) for X (or encodes free-text query with all-MiniLM-L6-v2).",
            "Neo4j native vector index returns top-50 candidates by cosine similarity.",
            "ReceptionScorer aggregates :PlatformReception nodes into per-platform sentiment.",
            "Optional cross-encoder (bge-reranker-base) or LightGBM LambdaRank fuses 6 features (sim, cf, reception, diversity, ce, freshness).",
            "Final ranked list returned with per-candidate breakdown — the same shape the dashboard renders.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Ingestion Data Flow (privacy-preserving)",
        "items": [
            "Scheduler triggers ingest_daily.py (cron or manual).",
            "Query :Book WHERE description AND thumbnail exist LIMIT N (skips 6 Sri Lankan seed books — no English-language chatter).",
            "For each (book × collector): call platform API, extract mention text.",
            "Batch sentiment: twitter-roberta-base-sentiment-latest (batch_size=16) → labels aggregated to (positive, neutral, negative).",
            "Persist only counts + external_ids; raw mention text is discarded before writing.",
            "MERGE (b)-[:RECEPTION_ON]->(r:PlatformReception) SET counts, ingested_at, expires_at (YouTube 30-day), demo=false.",
        ],
    },
    {
        "kind": "kv",
        "title": "Neo4j Data Model — Nodes",
        "rows": [
            (":Person", "id, username, displayName, bio, domain, publicKey, privateKey, passwordHash, avatarUrl, createdAt"),
            (":Book", "id, isbn, title, author, year, description, thumbnail, subjects, embedding (vector), goodbooksBookId, sourceCatalog"),
            (":Review", "id, content, rating, published, activityId, seedKey?"),
            (":Annotation", "id, motivation, bodyValue, exactText, prefix, suffix, startOffset, endOffset, bookSource, created, seedKey?"),
            (":PlatformReception", "platform, book_isbn, positive, neutral, negative, mentions, external_ids, ingested_at, expires_at, demo"),
        ],
    },
    {
        "kind": "kv",
        "title": "Neo4j Data Model — Relationships & Indexes",
        "rows": [
            ("Social edges", "(:Person)-[:FOLLOWS]->(:Person)  ·  (:Person)-[:LIKES]->(:Review)"),
            ("Authoring", "(:Person)-[:AUTHORED]->(:Review)  ·  (:Person)-[:ANNOTATED]->(:Annotation)"),
            ("Targets", "(:Review)-[:REVIEWS]->(:Book)  ·  (:Annotation)-[:ON_SOURCE]->(:Book)  ·  (:Book)-[:RECEPTION_ON]->(:PlatformReception)"),
            ("Uniqueness", "CONSTRAINT :Person(username) UNIQUE  ·  CONSTRAINT :Book(isbn) UNIQUE"),
            ("Vector index", "bookEmbedding on :Book(embedding) — cosine, 384 dims (MiniLM-L6-v2)"),
            ("Graph projection", "GDS projection over :Person → :Person for personalised PageRank"),
        ],
    },
    {
        "kind": "kv",
        "title": "Deployment Topology",
        "rows": [
            ("React FE", "native  ·  npm start  ·  :3000  ·  CRA proxy → :3001"),
            ("Express BE", "native  ·  nodemon  ·  :3001"),
            ("Neo4j 5.15", "docker  ·  :7474 (browser) + :7687 (bolt)  ·  volumes: neo4j_data, neo4j_logs"),
            ("Redis 7", "docker  ·  :6379"),
            ("PostgreSQL 16", "docker  ·  :5432  ·  volume: postgres_data"),
            ("ML service", "docker (built)  ·  :8000  ·  volumes: ml_hf_cache, ml_torch_cache, ./ml-service"),
            ("Dashboard", "docker (built)  ·  :8501  ·  same image as ml-service, streamlit entrypoint"),
        ],
    },
    {
        "kind": "kv",
        "title": "Recommender Strategies (all live)",
        "rows": RECOMMENDER_ROWS,
    },
    {
        "kind": "kv",
        "title": "Cross-Platform Ingestion",
        "rows": INGEST_ROWS,
    },
    {
        "kind": "bullets",
        "title": "Federation — ActivityPub",
        "items": [
            "Every user has an ActivityPub Person actor at /users/:username.",
            "WebFinger discovery at /.well-known/webfinger enables acct:user@domain lookup.",
            "HTTP Signatures (RSA-2048, node-forge) sign every outbound activity.",
            "Shared inbox at /inbox accepts Follow, Like, Announce, and Create activities.",
            "Reviews and annotations are exchanged as ActivityStreams objects with the W3C Anno JSON-LD context.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Annotations — W3C Web Annotation",
        "items": [
            "Full W3C Web Annotation Data Model with dual selectors — TextQuoteSelector + TextPositionSelector.",
            "Reader selects a passage; UI captures exact text, prefix/suffix, and character offsets.",
            "Annotation motivations supported: commenting, highlighting, tagging.",
            "Annotations propagate over ActivityPub as first-class Create activities.",
        ],
    },
    {
        "kind": "section",
        "kicker": "05",
        "title": "Delivery",
    },
    {
        "kind": "bullets",
        "title": "Timeline & Milestones",
        "items": [
            "Slice A — federation core + annotations + React FE  (done)",
            "Slice B — semantic recommender + LightFM CF + Neo4j vector index  (done)",
            "Slice C — cross-platform ingestion + reception scoring + LTR fusion  (done)",
            "Slice D — user study, evaluation, thesis writing  (current)",
        ],
    },
    {
        "kind": "bullets",
        "title": "Reproducibility",
        "items": [
            "One-command bring-up: docker compose up -d starts Neo4j, Postgres, Redis, ML service, dashboard.",
            "Idempotent seeders: seed.js (FE) and seed_kaggle_books.py (catalogue) both use MERGE, never wipe.",
            "Models are cached offline in the container (HF_HUB_OFFLINE=1) — no network needed after first build.",
            "All source, seeds, and configuration in one public git repository.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Limitations & Future Work",
        "items": [
            "LightFM is user-cold-start bound — needs onboarding taste-prompt to help new users.",
            "Cross-platform sentiment sample is small for less-discussed titles; needs longer accumulation window.",
            "Federation currently point-to-point; no relay support (BookWyrm-style relay planned).",
            "Cross-encoder rerank is slow (~500 MB model); could be replaced with smaller distilled reranker.",
        ],
    },
    {
        "kind": "quote",
        "title": "Thank You",
        "quote": "Questions?",
        "attribution": "IT22922670  ·  Supervised by Ms. Dinithi Pandithage",
    },
]

# ----- Technical walkthrough deck ------------------------------------------

TECHNICAL_SLIDES: list[dict[str, Any]] = [
    {
        "kind": "title",
        "title": "FedBook-Sem",
        "subtitle": "Technical Walkthrough — Architecture, Stack, Data Flow",
        "footer": "For engineers and technical reviewers",
    },
    {
        "kind": "bullets",
        "title": "What FedBook-Sem Is",
        "items": [
            SOLUTION_STATEMENT,
            "Two independent runtimes (Node backend, Python ML service) share one Neo4j.",
            "Streamlit dashboard is a thin client over the ML FastAPI.",
            "Idempotent seeders — FE seed and ML seed can co-exist safely.",
        ],
    },
    {
        "kind": "kv",
        "title": "Runtime Topology",
        "rows": [
            ("localhost:3000", "React FE (CRA dev server, proxies /api → :3001)"),
            ("localhost:3001", "Express backend (Node 20)"),
            ("localhost:7474 / 7687", "Neo4j 5.15 (browser + bolt)"),
            ("localhost:6379", "Redis 7"),
            ("localhost:5432", "PostgreSQL 16"),
            ("localhost:8000", "FastAPI ML service (Python 3.11)"),
            ("localhost:8501", "Streamlit dashboard"),
        ],
    },
    {
        "kind": "kv",
        "title": "Technology Stack",
        "rows": TECH_ROWS,
    },
    {
        "kind": "quote",
        "title": "System Diagram",
        "quote": (
            "Open docs/presentation/FedBookSem_Architecture.drawio in draw.io\n"
            "or app.diagrams.net for the layered architecture with all services,\n"
            "endpoints, and data flows."
        ),
        "attribution": "Client · Application · Ingestion · Data · External layers",
    },
    {
        "kind": "kv",
        "title": "Layered Architecture — Five Tiers",
        "rows": [
            ("Client", "React FE (:3000)  ·  Streamlit dashboard (:8501)  ·  remote ActivityPub servers"),
            ("Application", "Express backend (:3001)  ·  FastAPI ML service (:8000)  ·  ActivityPub inbox / outbox / WebFinger"),
            ("Ingestion", "Reddit / YouTube / Bluesky / Mastodon collectors → entity resolver → sentiment scorer → aggregator"),
            ("Data", "Neo4j 5.15 (:7687)  ·  Redis 7 (:6379)  ·  PostgreSQL 16 (:5432)  ·  cached model files on disk"),
            ("External", "Open Library  ·  Hardcover GraphQL  ·  HuggingFace Hub  ·  Kaggle 7k  ·  goodbooks-10k  ·  remote AP peers"),
        ],
    },
    {
        "kind": "bullets",
        "title": "Request Lifecycle — Login and Feed",
        "items": [
            "Browser POSTs /api/auth/login → CRA dev proxy (setupProxy.js) forwards to Express :3001.",
            "Express: MATCH (:Person {username}), bcrypt.compare(password, passwordHash), jwt.sign(actor, JWT_SECRET, '7d').",
            "JWT stored in localStorage; axios interceptor at src/api/client.js attaches Bearer to every subsequent call.",
            "GET /api/feed → verifyJwt middleware → single Cypher: MATCH (me)-[:FOLLOWS]->(f)-[:AUTHORED]->(r)-[:REVIEWS]->(b) RETURN ...",
            "Response JSON hydrated client-side, rendered by Feed page component.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Federation Flow — Outbound Follow",
        "items": [
            "POST /api/social/follow {targetId: 'acct:bob@remote.social'}.",
            "backend/src/routes/social.js → WebFinger lookup: GET https://remote.social/.well-known/webfinger?resource=acct:bob@remote.social.",
            "Fetch actor JSON from returned URL; extract publicKey and inbox.",
            "activitypub/delivery.js builds Follow activity, canonicalises headers, signs with alice.privateKey (RSA-SHA256).",
            "Signature header: keyId=<alice.publicKey.id>, algorithm=rsa-sha256, headers=(request-target) host date digest, signature=<b64>.",
            "POST to bob's inbox with Content-Type application/activity+json; expect 202 Accepted.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Federation Flow — Inbound Activity",
        "items": [
            "Remote server POSTs Follow / Like / Announce / Create to /inbox.",
            "backend/src/routes/inbox.js: parse Signature header, fetch remote actor's publicKey via keyId URL.",
            "Verify canonicalised headers against signature; reject on mismatch.",
            "Dispatch by activity type: Follow → MERGE (:FOLLOWS) + send Accept; Like → MERGE (:LIKES); Create Review → MERGE (:Review) + REVIEWS edge.",
            "All writes go through the same Neo4j driver used by the FE routes — one graph, many writers.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Recommender Request Flow — /recommend/similar",
        "items": [
            "GET /recommend/similar?isbn=X&rerank=ltr&limit=10.",
            "VectorSearcher: read :Book(embedding) for X (or encode ?text= with all-MiniLM-L6-v2).",
            "Neo4j: CALL db.index.vector.queryNodes('bookEmbedding', 50, $emb) YIELD node, score.",
            "ReceptionScorer: MATCH (b)-[:RECEPTION_ON]->(r) for candidate ISBNs → per-platform sentiment (Reddit .35, YT .25, BS/Ma .20).",
            "Cross-encoder (bge-reranker-base) computes ce_score for query × candidate pairs.",
            "LightGBM LambdaRank (ltr.pkl) fuses (sim, cf, reception, diversity, ce, freshness) in fixed FEATURE_ORDER.",
            "Sort by fused score, take top-N, optional Open Library / Hardcover enrichment layer.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Ingestion Data Flow (privacy-preserving)",
        "items": [
            "ingest_daily.py --platforms ... --limit N — cron or manual.",
            "MATCH (b:Book) WHERE b.description IS NOT NULL AND b.thumbnail IS NOT NULL LIMIT $limit.",
            "For each (book × collector): call platform API (PRAW / YouTube DataAPI / AT Protocol / Mastodon timeline).",
            "SentimentScorer.batch (twitter-roberta, batch=16) → labels aggregated via scorer.aggregate.",
            "Persist: MERGE (b)-[:RECEPTION_ON]->(r:PlatformReception {platform, book_isbn}) SET positive, neutral, negative, mentions, external_ids, ingested_at, expires_at (YouTube +30d), demo=false.",
            "Raw mention text never written; external_ids kept only for deletion sweeps.",
        ],
    },
    {
        "kind": "kv",
        "title": "Neo4j Data Model — Nodes",
        "rows": [
            (":Person", "id, username, displayName, bio, domain, publicKey, privateKey, passwordHash, avatarUrl, createdAt"),
            (":Book", "id, isbn, title, author, year, description, thumbnail, subjects, embedding (vector 384-d), goodbooksBookId, sourceCatalog"),
            (":Review", "id, content, rating (1-5), published, activityId, seedKey?"),
            (":Annotation", "id, motivation, bodyValue, exactText, prefix, suffix, startOffset, endOffset, bookSource, created, seedKey?"),
            (":PlatformReception", "platform, book_isbn, positive, neutral, negative, mentions, external_ids, ingested_at, expires_at, demo"),
        ],
    },
    {
        "kind": "kv",
        "title": "Neo4j Data Model — Relationships & Indexes",
        "rows": [
            ("Social", "(:Person)-[:FOLLOWS]->(:Person)  ·  (:Person)-[:LIKES]->(:Review)"),
            ("Authoring", "(:Person)-[:AUTHORED]->(:Review)  ·  (:Person)-[:ANNOTATED]->(:Annotation)"),
            ("Content targets", "(:Review)-[:REVIEWS]->(:Book)  ·  (:Annotation)-[:ON_SOURCE]->(:Book)"),
            ("Reception", "(:Book)-[:RECEPTION_ON]->(:PlatformReception)"),
            ("Constraints", "UNIQUE :Person(username)  ·  UNIQUE :Book(isbn)"),
            ("Vector index", "bookEmbedding on :Book(embedding), cosine, 384 dims (MiniLM-L6-v2)"),
            ("GDS", "projection :Person → :Person for personalised PageRank in /recommend/graph"),
        ],
    },
    {
        "kind": "kv",
        "title": "Deployment Topology",
        "rows": [
            ("React FE", "native  ·  react-scripts start  ·  :3000  ·  CRA proxy → :3001"),
            ("Express BE", "native  ·  nodemon src/app.js  ·  :3001"),
            ("Neo4j 5.15", "docker  ·  fedbooksem-neo4j  ·  :7474 + :7687  ·  APOC + GDS  ·  volumes: neo4j_data, neo4j_logs"),
            ("Redis 7", "docker  ·  fedbooksem-redis  ·  :6379"),
            ("PostgreSQL 16", "docker  ·  fedbooksem-postgres  ·  :5432  ·  volume: postgres_data"),
            ("ML service", "docker (built)  ·  fedbook-ml  ·  :8000  ·  volumes: ml_hf_cache, ml_torch_cache, ./ml-service"),
            ("Dashboard", "docker (built)  ·  fedbook-ml-dashboard  ·  :8501  ·  streamlit entrypoint"),
        ],
    },
    {
        "kind": "kv",
        "title": "Port Map (Host → Container / Process)",
        "rows": [
            ("3000", "React CRA dev server (native)"),
            ("3001", "Express backend (native)"),
            ("7474", "Neo4j HTTP browser UI (docker)"),
            ("7687", "Neo4j bolt (docker)"),
            ("6379", "Redis (docker)"),
            ("5432", "PostgreSQL (docker)"),
            ("8000", "FastAPI ML service (docker)"),
            ("8501", "Streamlit dashboard (docker)"),
        ],
    },
    {
        "kind": "kv",
        "title": "Model Registry (loaded at ML startup)",
        "rows": [
            ("all-MiniLM-L6-v2", "sentence-transformers  ·  semantic embeddings (384-d)  ·  ~90 MB  ·  ml_hf_cache"),
            ("bge-reranker-base", "BAAI cross-encoder  ·  query × candidate re-scoring  ·  ~500 MB  ·  ml_hf_cache"),
            ("twitter-roberta-base-sentiment-latest", "CardiffNLP RoBERTa  ·  sentiment for ingestion  ·  ~500 MB  ·  ml_hf_cache"),
            ("lightfm.pkl", "LightFM WARP  ·  collaborative filtering  ·  ~10 MB  ·  ml-service/models/  ·  trained off-machine"),
            ("ltr.pkl", "LightGBM LambdaRank  ·  learned rank fusion  ·  ~1 MB  ·  ml-service/models/"),
        ],
    },
    {
        "kind": "kv",
        "title": "Environment Configuration",
        "rows": [
            ("DOMAIN, BASE_URL", "Federation identity — used as ActivityPub actor IDs (e.g. http://localhost:3001)"),
            ("NEO4J_URI / USER / PASSWORD", "Shared by backend and ML — bolt://localhost:7687 native, bolt://neo4j:7687 in-docker"),
            ("JWT_SECRET", "Backend HS256 signing key — required or /api/auth/login crashes"),
            ("ML_SERVICE_URL", "Backend → ML gateway (http://ml-service:8000 in docker network)"),
            ("YOUTUBE_API_KEY / REDDIT_CLIENT_ID+SECRET / BSKY_HANDLE+APP_PASSWORD", "Ingestion auth (Bluesky and Mastodon work anonymously)"),
            ("HARDCOVER_API_TOKEN", "Optional — live star ratings in book cards"),
            ("MASTODON_INSTANCES", "Comma-separated hashtag-timeline hosts (default: mastodon.social, ohai.social)"),
            ("HF_HUB_OFFLINE, TRANSFORMERS_OFFLINE", "Force cached HuggingFace weights — no network at startup"),
        ],
    },
    {
        "kind": "bullets",
        "title": "Security & Signature Design",
        "items": [
            "Signing keys — RSA-2048 pair generated at user creation by node-forge, stored on the :Person node (publicKey + privateKey).",
            "Outbound activity — HTTP Signatures over (request-target) host date digest; algorithm=rsa-sha256; keyId links to publicKey URL.",
            "Inbound activity — /inbox fetches remote actor's publicKey via keyId, verifies canonicalised headers, rejects on mismatch.",
            "JWT — HS256 with JWT_SECRET, 7-day expiry; payload = {id, username, displayName, bio, domain, avatarUrl}.",
            "Passwords — bcryptjs, cost factor 12; comparison in constant time.",
            "Helmet middleware sets CSP, HSTS, X-Frame-Options, etc.; response headers confirm.",
        ],
    },
    {
        "kind": "kv",
        "title": "Data Contracts (JSON-LD / JSON shapes)",
        "rows": [
            ("Actor (AP)", '{@context, type: "Person", id, preferredUsername, inbox, outbox, followers, following, publicKey: {id, owner, publicKeyPem}}'),
            ("Follow (AP)", '{@context, type: "Follow", actor, object: <target actor URL>, id}'),
            ("Web Annotation", '{@context, type: "Annotation", motivation, body: {type: "TextualBody", value}, target: {source, selector: [TextQuoteSelector, TextPositionSelector]}}'),
            ("PlatformReception", '{platform, book_isbn, positive, neutral, negative, mentions, external_ids, ingested_at, expires_at, demo}'),
            ("JWT payload", '{id, username, displayName, bio, domain, avatarUrl, iat, exp}'),
        ],
    },
    {
        "kind": "bullets",
        "title": "Express Backend — Route Groups",
        "items": [
            "/api/auth — POST /login → JWT (7d) signed with process.env.JWT_SECRET.",
            "/api/reviews — CRUD + like + announce + replies; federates via ActivityPub CREATE.",
            "/api/annotations — W3C Web Annotation CRUD with TextQuote + TextPosition selectors.",
            "/api/feed — personalised, all, recommendations, liked-by-followed.",
            "/api/books, /api/users, /api/social, /api/covers (SVG cover generator).",
            "/users/:u, /inbox, /.well-known/webfinger — ActivityPub actor + WebFinger surface.",
        ],
    },
    {
        "kind": "bullets",
        "title": "React Frontend — Composition",
        "items": [
            "Create React App with a custom setupProxy.js forwarding /api, /avatars, /users, /books, /.well-known, /inbox → :3001.",
            "Single axios client (src/api/client.js) with a JWT bearer interceptor.",
            "React Router v6 with pages: Feed, Book, BookReader, Profile, Login, Register, Followers, Following, User.",
            "Custom design system — Playfair Display + DM Sans over CSS variables.",
        ],
    },
    {
        "kind": "kv",
        "title": "ML Service — FastAPI Endpoints",
        "rows": [
            ("GET /health", "Model load state (transformer, cross-encoder, LightFM, LTR)"),
            ("GET /books/search", "Fuzzy title / author lookup"),
            ("GET /books/details", "Catalogue + reception + optional live enrichment"),
            ("GET /recommend/similar", "Semantic k-NN with optional cross-encoder / linear / LTR re-rank"),
            ("GET /recommend/cf", "LightFM WARP — 503 if lightfm.pkl missing"),
            ("GET /recommend/graph", "Neo4j GDS k-NN and personalised PageRank"),
        ],
    },
    {
        "kind": "kv",
        "title": "Recommender Strategies (all live)",
        "rows": RECOMMENDER_ROWS,
    },
    {
        "kind": "bullets",
        "title": "Sentiment Pipeline",
        "items": [
            "CardiffNLP twitter-roberta-base-sentiment-latest — batch=16, single or batch inference.",
            "Twitter-style preprocessing: @username → @user, URLs → http.",
            "Offline fallback lexicon scorer activates when HF hub or model file is unavailable — labels neutral by default, requires strong signal to shift.",
            "Aggregation returns (positive, neutral, negative) counts per (book, platform).",
        ],
    },
    {
        "kind": "bullets",
        "title": "Entity Resolution",
        "items": [
            "rapidfuzz WRatio fuzzy match on title + author, default min_score = 82.",
            "In-memory catalogue index built at startup, refreshed periodically.",
            "Maps user-generated mention text (misspelled, no ISBN) to canonical :Book.isbn before scoring.",
        ],
    },
    {
        "kind": "kv",
        "title": "Cross-Platform Ingestion",
        "rows": INGEST_ROWS,
    },
    {
        "kind": "bullets",
        "title": "Storage Model — Neo4j Labels & Relationships",
        "items": [
            "Nodes: :Person, :Book, :Review, :Annotation, :PlatformReception.",
            "Relationships: :FOLLOWS, :AUTHORED, :REVIEWS, :LIKES, :ANNOTATED, :ON_SOURCE, :RECEPTION_ON.",
            "Native vector index on :Book(embedding) for semantic k-NN.",
            "GDS graph projection over :Person → :Person for personalised PageRank.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Federation Details",
        "items": [
            "ActivityPub actor JSON-LD at /users/:username with publicKey block for HTTP Signatures.",
            "WebFinger at /.well-known/webfinger resolves acct:user@domain to actor URL.",
            "Outbound activities signed with node-forge (RSA-2048, private key stored on :Person node).",
            "Shared inbox at /inbox verifies signature, dispatches Follow / Like / Announce / Create.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Deployment — Docker Compose",
        "items": [
            "Services: neo4j, redis, postgres, ml-service, ml-dashboard.",
            "ml-service and ml-dashboard both build from ml-service/Dockerfile; dashboard runs streamlit.",
            "Volumes: neo4j_data, postgres_data, ml_hf_cache, ml_torch_cache — model + DB state persists across rebuilds.",
            "Backend + frontend run natively via nodemon / react-scripts for hot reload.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Operations",
        "items": [
            "npm run seed — idempotent FE seed (3 users, 6 Sri Lankan books, follows, reviews, annotations).",
            "seed_kaggle_books.py — Kaggle 7k catalogue via MERGE (no wipe).",
            "build_embeddings.py — MiniLM embeddings written to :Book(embedding).",
            "ingest_daily.py --platforms ... --limit N — real cross-platform ingestion, --drop-mock replaces demo data.",
            "purge_expired_receptions.py — daily YouTube 30-day retention purge.",
            "reddit_deletion_sweep.py — weekly, honours Reddit deletions.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Privacy & Compliance Design",
        "items": [
            "Raw mention text is never persisted — only aggregated counts and external IDs for deletion sweep.",
            "YouTube receptions carry expires_at, purged after 30 days.",
            "Reddit deletion sweep removes counts for posts that no longer exist.",
            "All external calls carry a project user-agent identifying the researcher.",
        ],
    },
    {
        "kind": "quote",
        "title": "Where to Look in the Repo",
        "quote": (
            "backend/src/routes/*.js  ·  backend/src/activitypub/delivery.js\n"
            "ml-service/src/fedbook_ml/api.py  ·  ml-service/src/fedbook_ml/ingest/*.py\n"
            "ml-service/src/fedbook_ml/dashboard/pages/*.py"
        ),
        "attribution": "docs/SETUP_ML_AND_DASHBOARD.md walks the whole bring-up.",
    },
]

# ----- General overview deck -----------------------------------------------

GENERAL_SLIDES: list[dict[str, Any]] = [
    {
        "kind": "title",
        "title": "FedBook-Sem",
        "subtitle": "A social bookstore that talks to other social bookstores",
        "footer": "Overview  ·  IT22922670",
    },
    {
        "kind": "bullets",
        "title": "What is FedBook-Sem?",
        "items": [
            "An online bookstore with a social layer — reviews, annotations, followers.",
            "But unlike Goodreads or Storygraph, it is federated — instances can talk to each other and to Mastodon.",
            "It also learns from the wider web — reading what people say about books on Reddit, YouTube, Bluesky, and Mastodon — to make better recommendations.",
        ],
    },
    {
        "kind": "two_col",
        "title": "Why It Matters",
        "left": {
            "heading": "Today",
            "items": [
                "Your reviews live inside one product.",
                "Your followers cannot leave with you.",
                "Recommendations use only the platform's own data.",
                "Annotations are locked to one app.",
            ],
        },
        "right": {
            "heading": "With FedBook-Sem",
            "items": [
                "Your identity is a portable ActivityPub actor.",
                "Followers reach you across servers.",
                "Recommendations blend semantic, social, and cross-platform signal.",
                "Annotations follow a W3C standard used in academic research.",
            ],
        },
    },
    {
        "kind": "bullets",
        "title": "What You Can Do",
        "items": [
            "Browse a curated book catalogue with covers, subjects, and star ratings.",
            "Write reviews, like and reshare reviews from other users.",
            "Annotate passages inside the book reader — highlight, comment, or tag.",
            "Follow other readers and see their activity in a personalised feed.",
            "Get recommendations that also consider what people are saying about a book on social media.",
        ],
    },
    {
        "kind": "kv",
        "title": "How the Pieces Fit",
        "rows": [
            ("React app", "The interface you use in your browser."),
            ("Express API", "The backend that handles reviews, follows, and federation."),
            ("Neo4j graph", "Stores the social network — who follows whom, what they read."),
            ("ML service", "Ranks and recommends books using several models."),
            ("Streamlit dashboard", "A developer view that shows how the recommender scored each book."),
            ("Ingestion worker", "Pulls book chatter from Reddit, YouTube, Bluesky, Mastodon (privately)."),
        ],
    },
    {
        "kind": "kv",
        "title": "The Architecture in Layers",
        "rows": [
            ("What you see", "The React app in your browser and the Streamlit dashboard for developers."),
            ("What handles requests", "Two servers — a Node.js one for the social features, a Python one for the recommender."),
            ("Where knowledge lives", "A Neo4j graph database — think of it as a giant map of people, books, and their connections."),
            ("Supporting services", "Redis for fast caching, PostgreSQL as a backup relational store."),
            ("What it talks to outside", "Other ActivityPub servers (Mastodon-style), plus Reddit / YouTube / Bluesky / Mastodon for reception signal."),
        ],
    },
    {
        "kind": "bullets",
        "title": "What Happens When You Follow Someone",
        "items": [
            "You click Follow on a user profile.",
            "Your server looks up the other person via WebFinger (like an email lookup for social identities).",
            "Your server sends a cryptographically signed 'Follow' message to their server's inbox.",
            "Their server checks the signature is really yours, records the follow, replies with an 'Accept'.",
            "You now see each other's activity — even if you're on different servers.",
        ],
    },
    {
        "kind": "bullets",
        "title": "What Happens When You Open a Book Page",
        "items": [
            "The React app asks the ML service for similar books.",
            "ML converts the book into a mathematical fingerprint (an embedding vector).",
            "It finds the 50 most similar fingerprints already in the database.",
            "Re-scores them using: what people are saying about them online, who else read them, how fresh the chatter is.",
            "Returns the top 10 with a mini scorecard for each — sentiment per platform, star rating, subjects.",
        ],
    },
    {
        "kind": "kv",
        "title": "Where Different Kinds of Data Live",
        "rows": [
            ("Neo4j graph database", "People, books, follows, reviews, annotations, per-platform reception counts."),
            ("Redis", "Fast temporary data — session cache, background job queues."),
            ("PostgreSQL", "Reserved for future relational data (e.g. purchase history)."),
            ("Disk (inside the ML container)", "Downloaded AI models — about 1 GB total, cached so no internet is needed at startup."),
            ("Nowhere", "The actual text of anyone's Reddit / YouTube / Bluesky / Mastodon posts — deliberately not stored."),
        ],
    },
    {
        "kind": "bullets",
        "title": "How Recommendations Work — In Plain Terms",
        "items": [
            "Semantic similarity — 'books that read like this one'.",
            "Collaborative filtering — 'readers like you also liked'.",
            "Cross-platform reception — 'what is the internet saying about this book right now?'.",
            "A learned model blends these signals into a single ranked list.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Privacy by Design",
        "items": [
            "The system never stores the actual text of anyone's Reddit / YouTube / Bluesky / Mastodon posts.",
            "It only keeps counts of how many mentions were positive, neutral, or negative.",
            "YouTube data is automatically deleted after 30 days.",
            "Reddit deletions are honoured — if the original post is gone, our count is scrubbed.",
        ],
    },
    {
        "kind": "bullets",
        "title": "Built On Open Standards",
        "items": [
            "ActivityPub — the same protocol Mastodon uses. Any ActivityPub server can talk to FedBook-Sem.",
            "W3C Web Annotation — the same standard used in scholarly text and hypothes.is.",
            "WebFinger — lets you find users as acct:name@domain, just like email.",
            "HTTP Signatures — every message between servers is cryptographically signed.",
        ],
    },
    {
        "kind": "kv",
        "title": "Under the Hood",
        "rows": [
            ("Backend", "Node.js + Express + Neo4j"),
            ("Frontend", "React + custom design"),
            ("ML", "Python + FastAPI + Hugging Face + LightFM"),
            ("Dashboard", "Streamlit"),
            ("Infrastructure", "Docker Compose (Neo4j, Redis, PostgreSQL)"),
        ],
    },
    {
        "kind": "bullets",
        "title": "How It Runs",
        "items": [
            "docker compose up -d brings up the databases and the ML services.",
            "npm run dev:backend and npm run dev:frontend start the app locally.",
            "Sign in as one of the demo users (alice / bob / carol) and start reading, reviewing, annotating.",
            "The dashboard at :8501 lets you inspect the recommender's decisions.",
        ],
    },
    {
        "kind": "bullets",
        "title": "What's Next",
        "items": [
            "User study with real readers (task-completion, subjective satisfaction).",
            "Federation interop tests with Mastodon and BookWyrm.",
            "Faster reranker to trim /recommend/similar latency.",
            "Onboarding taste-prompt to help brand-new users get good recommendations.",
        ],
    },
    {
        "kind": "quote",
        "title": "Thank You",
        "quote": "Questions?",
        "attribution": "SLIIT R26-IT-107  ·  IT22922670",
    },
]

DECKS = {
    "FedBookSem_Academic_Viva": {
        "title": "FedBook-Sem — Academic Viva",
        "audience": "Academic viva / evaluation panel",
        "slides": ACADEMIC_SLIDES,
    },
    "FedBookSem_Technical_Walkthrough": {
        "title": "FedBook-Sem — Technical Walkthrough",
        "audience": "Technical stakeholder walkthrough",
        "slides": TECHNICAL_SLIDES,
    },
    "FedBookSem_General_Overview": {
        "title": "FedBook-Sem — Overview",
        "audience": "General overview",
        "slides": GENERAL_SLIDES,
    },
}


# ---------------------------------------------------------------------------
# PPTX renderer
# ---------------------------------------------------------------------------


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _add_rect(slide, x_in, y_in, w_in, h_in, fill_hex, line=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    if not line:
        shape.line.fill.background()
    return shape


def _add_text(slide, x_in, y_in, w_in, h_in, text, *, size=18, bold=False,
              color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = _hex_to_rgb(color)
    return tb


def _add_bullets(slide, x_in, y_in, w_in, h_in, items, *, size=16, color=INK,
                 font="Calibri", spacing=6, bullet_char="•"):
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(size)
        run.font.name = font
        run.font.color.rgb = _hex_to_rgb(color)
    return tb


def _footer_bar(slide, deck_title):
    # Thin bottom bar with deck title on left, page number placeholder on right.
    _add_rect(slide, 0, SLIDE_H_IN - 0.35, SLIDE_W_IN, 0.35, PRIMARY)
    _add_text(slide, 0.4, SLIDE_H_IN - 0.32, 8, 0.3, deck_title,
              size=10, color="#FFFFFF")


def _slide_title(slide, title, deck_title):
    # Accent bar on left, title text, primary underline
    _add_rect(slide, 0.4, 0.5, 0.12, 0.55, ACCENT)
    _add_text(slide, 0.7, 0.4, 12, 0.8, title, size=28, bold=True, color=PRIMARY)
    _add_rect(slide, 0.4, 1.15, 12.4, 0.03, PRIMARY)
    _footer_bar(slide, deck_title)


def render_title(slide, s, deck_title):
    _add_rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, PRIMARY)
    _add_rect(slide, 0.6, 2.9, 0.18, 1.7, ACCENT)
    _add_text(slide, 0.9, 2.4, 12, 1.4, s["title"], size=54, bold=True,
              color="#FFFFFF")
    _add_text(slide, 0.9, 3.9, 12, 0.9, s.get("subtitle", ""), size=22,
              color="#E7EAF0")
    _add_text(slide, 0.9, 5.7, 12, 0.5, s.get("footer", ""), size=13,
              color=ACCENT)


def render_section(slide, s, deck_title):
    _add_rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, PRIMARY)
    if s.get("kicker"):
        _add_text(slide, 0.8, 2.6, 12, 0.6, s["kicker"], size=20, bold=True,
                  color=ACCENT, font="Calibri")
    _add_text(slide, 0.8, 3.1, 12, 1.4, s["title"], size=46, bold=True,
              color="#FFFFFF")
    _add_rect(slide, 0.85, 4.5, 4.5, 0.05, ACCENT)


def render_bullets(slide, s, deck_title):
    _slide_title(slide, s["title"], deck_title)
    _add_bullets(slide, 0.7, 1.55, 12, 5.2, s["items"], size=17)
    if s.get("note"):
        _add_text(slide, 0.7, 6.75, 12, 0.35, s["note"], size=11, color=MUTED)


def render_two_col(slide, s, deck_title):
    _slide_title(slide, s["title"], deck_title)
    col_w = 5.9
    left, right = s["left"], s["right"]
    _add_text(slide, 0.7, 1.5, col_w, 0.5, left["heading"], size=16,
              bold=True, color=PRIMARY)
    _add_bullets(slide, 0.7, 2.0, col_w, 4.8, left["items"], size=15)
    _add_text(slide, 7.0, 1.5, col_w, 0.5, right["heading"], size=16,
              bold=True, color=PRIMARY)
    _add_bullets(slide, 7.0, 2.0, col_w, 4.8, right["items"], size=15)


def render_kv(slide, s, deck_title):
    _slide_title(slide, s["title"], deck_title)
    rows = s["rows"]
    top = 1.6
    row_h = min(0.55, (5.3 / max(len(rows), 1)))
    for i, (k, v) in enumerate(rows):
        y = top + i * row_h
        _add_rect(slide, 0.7, y + 0.02, 12.2, row_h - 0.05,
                  SOFT if i % 2 == 0 else "#FFFFFF")
        _add_text(slide, 0.85, y + 0.05, 3.5, row_h - 0.1, k,
                  size=13, bold=True, color=PRIMARY)
        _add_text(slide, 4.4, y + 0.05, 8.4, row_h - 0.1, v,
                  size=12, color=INK)


def render_quote(slide, s, deck_title):
    _slide_title(slide, s["title"], deck_title)
    _add_rect(slide, 1.0, 2.2, 0.12, 3.0, ACCENT)
    _add_text(slide, 1.4, 2.1, 11, 3.2, s["quote"], size=22, color=INK)
    if s.get("attribution"):
        _add_text(slide, 1.4, 5.6, 11, 0.5, "— " + s["attribution"],
                  size=13, color=MUTED)


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "bullets": render_bullets,
    "two_col": render_two_col,
    "kv": render_kv,
    "quote": render_quote,
}


def build_pptx(out_path: Path, deck: dict[str, Any]):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    for s in deck["slides"]:
        slide = _blank_slide(prs)
        RENDERERS[s["kind"]](slide, s, deck["title"])
    prs.save(str(out_path))


# ---------------------------------------------------------------------------
# PDF renderer (reportlab)
# ---------------------------------------------------------------------------

PAGE_SIZE = landscape(A4)  # 297 × 210 mm
PAGE_W, PAGE_H = PAGE_SIZE


def _pdf_footer(c: rl_canvas.Canvas, deck_title: str, page_num: int):
    c.setFillColor(_hex_to_rl(PRIMARY))
    c.rect(0, 0, PAGE_W, 10 * mm, stroke=0, fill=1)
    c.setFillColor(rl_colors.white)
    c.setFont("Helvetica", 8)
    c.drawString(12 * mm, 3.5 * mm, deck_title)
    c.drawRightString(PAGE_W - 12 * mm, 3.5 * mm, f"{page_num}")


def _pdf_wrap(text, style, avail_w):
    """Yield lines from Paragraph.wrap for hand-drawing."""
    p = Paragraph(text.replace("\n", "<br/>"), style)
    _, h = p.wrap(avail_w, PAGE_H)
    return p, h


def _draw_paragraph(c, text, x, y, w, style):
    p = Paragraph(text.replace("\n", "<br/>"), style)
    _, h = p.wrap(w, PAGE_H)
    p.drawOn(c, x, y - h)
    return h


def _pdf_title_slide(c, s, deck_title, page_num):
    c.setFillColor(_hex_to_rl(PRIMARY))
    c.rect(0, 0, PAGE_W, PAGE_H, stroke=0, fill=1)
    c.setFillColor(_hex_to_rl(ACCENT))
    c.rect(20 * mm, PAGE_H / 2 - 20 * mm, 5 * mm, 40 * mm, stroke=0, fill=1)

    c.setFillColor(rl_colors.white)
    c.setFont("Helvetica-Bold", 34)
    c.drawString(32 * mm, PAGE_H / 2 + 8 * mm, s["title"])
    c.setFont("Helvetica", 16)
    c.setFillColor(rl_colors.Color(0.91, 0.92, 0.94))
    # subtitle can be long — wrap manually
    subtitle = s.get("subtitle", "")
    style = ParagraphStyle("sub", fontName="Helvetica", fontSize=15,
                           textColor=rl_colors.Color(0.91, 0.92, 0.94),
                           leading=19)
    p = Paragraph(subtitle, style)
    _, h = p.wrap(PAGE_W - 60 * mm, 30 * mm)
    p.drawOn(c, 32 * mm, PAGE_H / 2 - 6 * mm - h)
    c.setFillColor(_hex_to_rl(ACCENT))
    c.setFont("Helvetica-Oblique", 10)
    c.drawString(32 * mm, 25 * mm, s.get("footer", ""))


def _pdf_section_slide(c, s, deck_title, page_num):
    c.setFillColor(_hex_to_rl(PRIMARY))
    c.rect(0, 0, PAGE_W, PAGE_H, stroke=0, fill=1)
    if s.get("kicker"):
        c.setFillColor(_hex_to_rl(ACCENT))
        c.setFont("Helvetica-Bold", 14)
        c.drawString(30 * mm, PAGE_H / 2 + 20 * mm, s["kicker"])
    c.setFillColor(rl_colors.white)
    c.setFont("Helvetica-Bold", 30)
    c.drawString(30 * mm, PAGE_H / 2 + 5 * mm, s["title"])
    c.setFillColor(_hex_to_rl(ACCENT))
    c.rect(30 * mm, PAGE_H / 2 - 3 * mm, 45 * mm, 1.2 * mm, stroke=0, fill=1)


def _pdf_content_header(c, title, deck_title, page_num):
    # Accent + title bar
    c.setFillColor(_hex_to_rl(ACCENT))
    c.rect(15 * mm, PAGE_H - 22 * mm, 3 * mm, 10 * mm, stroke=0, fill=1)
    c.setFillColor(_hex_to_rl(PRIMARY))
    c.setFont("Helvetica-Bold", 18)
    c.drawString(22 * mm, PAGE_H - 20 * mm, title)
    c.setStrokeColor(_hex_to_rl(PRIMARY))
    c.setLineWidth(0.6)
    c.line(15 * mm, PAGE_H - 25 * mm, PAGE_W - 15 * mm, PAGE_H - 25 * mm)


def _pdf_bullets_slide(c, s, deck_title, page_num):
    _pdf_content_header(c, s["title"], deck_title, page_num)
    style = ParagraphStyle("bul", fontName="Helvetica", fontSize=12,
                           textColor=_hex_to_rl(INK), leading=17,
                           bulletIndent=0, leftIndent=6 * mm)
    y = PAGE_H - 32 * mm
    for item in s["items"]:
        text = f"•&nbsp;&nbsp;{item}"
        p = Paragraph(text, style)
        _, h = p.wrap(PAGE_W - 30 * mm, PAGE_H)
        p.drawOn(c, 15 * mm, y - h)
        y -= (h + 4 * mm)


def _pdf_two_col_slide(c, s, deck_title, page_num):
    _pdf_content_header(c, s["title"], deck_title, page_num)
    col_w = (PAGE_W - 45 * mm) / 2
    hstyle = ParagraphStyle("h", fontName="Helvetica-Bold", fontSize=13,
                            textColor=_hex_to_rl(PRIMARY), leading=17)
    bstyle = ParagraphStyle("b", fontName="Helvetica", fontSize=11,
                            textColor=_hex_to_rl(INK), leading=15)

    def _column(x, col):
        y = PAGE_H - 32 * mm
        p = Paragraph(col["heading"], hstyle)
        _, h = p.wrap(col_w, PAGE_H)
        p.drawOn(c, x, y - h)
        y -= (h + 3 * mm)
        for item in col["items"]:
            p = Paragraph(f"•&nbsp;&nbsp;{item}", bstyle)
            _, h = p.wrap(col_w, PAGE_H)
            p.drawOn(c, x, y - h)
            y -= (h + 3 * mm)

    _column(15 * mm, s["left"])
    _column(15 * mm + col_w + 15 * mm, s["right"])


def _pdf_kv_slide(c, s, deck_title, page_num):
    _pdf_content_header(c, s["title"], deck_title, page_num)
    rows = s["rows"]
    top = PAGE_H - 30 * mm
    max_h = top - 20 * mm
    row_h = min(11 * mm, max_h / max(len(rows), 1))
    key_style = ParagraphStyle("k", fontName="Helvetica-Bold", fontSize=10,
                               textColor=_hex_to_rl(PRIMARY), leading=13)
    val_style = ParagraphStyle("v", fontName="Helvetica", fontSize=10,
                               textColor=_hex_to_rl(INK), leading=13)
    for i, (k, v) in enumerate(rows):
        y = top - (i + 1) * row_h
        if i % 2 == 0:
            c.setFillColor(_hex_to_rl(SOFT))
            c.rect(15 * mm, y, PAGE_W - 30 * mm, row_h, stroke=0, fill=1)
        pk = Paragraph(k, key_style)
        pk.wrap(60 * mm, row_h)
        pk.drawOn(c, 18 * mm, y + row_h - 8 * mm)
        pv = Paragraph(v, val_style)
        pv.wrap(PAGE_W - 100 * mm, row_h)
        pv.drawOn(c, 82 * mm, y + row_h - 8 * mm)


def _pdf_quote_slide(c, s, deck_title, page_num):
    _pdf_content_header(c, s["title"], deck_title, page_num)
    c.setFillColor(_hex_to_rl(ACCENT))
    c.rect(28 * mm, PAGE_H / 2 - 25 * mm, 2 * mm, 50 * mm, stroke=0, fill=1)
    style = ParagraphStyle("q", fontName="Helvetica", fontSize=15,
                           textColor=_hex_to_rl(INK), leading=22)
    p = Paragraph(s["quote"].replace("\n", "<br/>"), style)
    _, h = p.wrap(PAGE_W - 60 * mm, 60 * mm)
    p.drawOn(c, 34 * mm, PAGE_H / 2 - 5 * mm)
    if s.get("attribution"):
        astyle = ParagraphStyle("a", fontName="Helvetica-Oblique", fontSize=10,
                                textColor=_hex_to_rl(MUTED))
        pa = Paragraph("— " + s["attribution"], astyle)
        pa.wrap(PAGE_W - 60 * mm, 20 * mm)
        pa.drawOn(c, 34 * mm, PAGE_H / 2 - 25 * mm)


PDF_RENDERERS = {
    "title": _pdf_title_slide,
    "section": _pdf_section_slide,
    "bullets": _pdf_bullets_slide,
    "two_col": _pdf_two_col_slide,
    "kv": _pdf_kv_slide,
    "quote": _pdf_quote_slide,
}


def build_pdf(out_path: Path, deck: dict[str, Any]):
    c = rl_canvas.Canvas(str(out_path), pagesize=PAGE_SIZE)
    for i, s in enumerate(deck["slides"], start=1):
        PDF_RENDERERS[s["kind"]](c, s, deck["title"], i)
        if s["kind"] not in ("title", "section"):
            _pdf_footer(c, deck["title"], i)
        c.showPage()
    c.save()


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------


def main():
    out_dir = Path(__file__).parent
    out_dir.mkdir(parents=True, exist_ok=True)
    for name, deck in DECKS.items():
        pptx_path = out_dir / f"{name}.pptx"
        pdf_path = out_dir / f"{name}.pdf"
        print(f"Building {name} ...")
        build_pptx(pptx_path, deck)
        build_pdf(pdf_path, deck)
        print(f"  wrote {pptx_path.name}, {pdf_path.name}  "
              f"({len(deck['slides'])} slides)")


if __name__ == "__main__":
    main()
